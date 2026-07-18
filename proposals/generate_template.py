"""
Cognitidata Proposal Template Generator
Produces: cognitidata-proposal-template.pptx
Design: editorial minimalist — mirrors the Cognitidata website aesthetic
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour tokens ────────────────────────────────────────────────────────────
BG      = RGBColor(0xF0, 0xF0, 0xF0)   # page background
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)   # card surface
TEXT    = RGBColor(0x11, 0x11, 0x11)   # primary text
BODY    = RGBColor(0x44, 0x44, 0x44)   # body copy
MUTED   = RGBColor(0x77, 0x77, 0x77)   # secondary body
HINT    = RGBColor(0xAA, 0xAA, 0xAA)   # eyebrow / hint
BORDER  = RGBColor(0xDC, 0xDC, 0xDC)   # rule / border
BTN_BG  = RGBColor(0x11, 0x11, 0x11)   # button / dark block
BTN_TXT = RGBColor(0xF0, 0xF0, 0xF0)   # button text

# ── Fonts ─────────────────────────────────────────────────────────────────────
SERIF = "Cormorant Garamond"   # headlines
SANS  = "Outfit"               # body / UI

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

MARGIN_L = Inches(1.1)
MARGIN_R = Inches(1.1)
MARGIN_T = Inches(0.85)
USABLE_W = W - MARGIN_L - MARGIN_R


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ── Primitive helpers ─────────────────────────────────────────────────────────

def fill_bg(slide, color=BG):
    """Set slide background to a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width_pt=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_rule(slide, y, color=BORDER, thickness_pt=0.75):
    """Draw a full-width horizontal rule."""
    return add_rect(slide, MARGIN_L, y, USABLE_W, Pt(thickness_pt),
                    fill_color=color, line_color=None)


def add_textbox(slide, x, y, w, h):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    return txb


def para_run(tf, text, font_name, font_size_pt, bold=False, italic=False,
             color=TEXT, align=PP_ALIGN.LEFT, space_before_pt=0, space_after_pt=0,
             tracking=0):
    """Add a paragraph with a single run to a text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before_pt)
    p.space_after  = Pt(space_after_pt)
    run = p.add_run()
    run.text = text
    f = run.font
    f.name  = font_name
    f.size  = Pt(font_size_pt)
    f.bold  = bold
    f.italic = italic
    f.color.rgb = color
    # Letter-spacing is not natively supported in python-pptx without XML hacking;
    # we skip it — the font choice alone carries the identity.
    return p


def eyebrow(tf, text, space_before_pt=0):
    return para_run(tf, text.upper(), SANS, 9, color=HINT,
                    space_before_pt=space_before_pt, space_after_pt=6)


def headline(tf, text, size=44, italic_wrap=False, space_before_pt=8, space_after_pt=0):
    """Main serif headline. Use *asterisks* to italicise portions — simplified: full italic flag."""
    return para_run(tf, text, SERIF, size, italic=italic_wrap, color=TEXT,
                    space_before_pt=space_before_pt, space_after_pt=space_after_pt)


def body_para(tf, text, size=13, color=BODY, space_before_pt=4, space_after_pt=0):
    return para_run(tf, text, SANS, size, color=color,
                    space_before_pt=space_before_pt, space_after_pt=space_after_pt)


def ghost_number(tf, number_str, size=80):
    """Large faint section number in the design-token border colour."""
    return para_run(tf, number_str, SERIF, size, color=BORDER, space_before_pt=0)


# ── Reusable layout blocks ────────────────────────────────────────────────────

def section_header(slide, eyebrow_text, title_text, title_size=44, y_start=None):
    """
    Renders the standard section header:
      EYEBROW LABEL
      Large Serif Title
    Returns the bottom y of the block.
    """
    y = y_start if y_start is not None else MARGIN_T
    txb = add_textbox(slide, MARGIN_L, y, USABLE_W, Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    # clear default empty paragraph
    tf.paragraphs[0].text = ""
    eyebrow(tf, eyebrow_text)
    headline(tf, title_text, size=title_size, space_before_pt=4)
    return y + Inches(1.5)


def numbered_item_row(slide, number_str, title, body, y, col_number_x=None,
                      col_content_x=None, row_height=Inches(1.1)):
    """
    Renders one numbered item:
      [01]   Title
             Body paragraph
    Returns bottom y of the row.
    """
    num_x = col_number_x if col_number_x is not None else MARGIN_L
    con_x = col_content_x if col_content_x is not None else MARGIN_L + Inches(0.85)
    num_w = Inches(0.75)
    con_w = USABLE_W - Inches(0.85)

    # Number
    txb_n = add_textbox(slide, num_x, y, num_w, row_height)
    tf_n  = txb_n.text_frame
    tf_n.paragraphs[0].text = ""
    para_run(tf_n, number_str, SERIF, 13, color=HINT, space_before_pt=2)

    # Content
    txb_c = add_textbox(slide, con_x, y, con_w, row_height)
    tf_c  = txb_c.text_frame
    tf_c.paragraphs[0].text = ""
    para_run(tf_c, title, SERIF, 18, color=TEXT, space_before_pt=0, space_after_pt=3)
    body_para(tf_c, body, size=12, color=MUTED, space_before_pt=2)

    return y + row_height


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_bg(slide)

    # Top nav bar
    nav_h = Inches(0.65)
    add_rect(slide, 0, 0, W, nav_h, fill_color=BG)
    add_rect(slide, 0, nav_h - Pt(0.75), W, Pt(0.75), fill_color=BORDER)

    # Logo wordmark
    txb = add_textbox(slide, MARGIN_L, Inches(0.16), Inches(3), Inches(0.4))
    tf = txb.text_frame
    tf.paragraphs[0].text = ""
    para_run(tf, "COGNITIDATA", SANS, 10, color=TEXT)

    # Prepared for label (right side of nav)
    txb2 = add_textbox(slide, W - MARGIN_R - Inches(3), Inches(0.16), Inches(3), Inches(0.4))
    tf2 = txb2.text_frame
    tf2.paragraphs[0].text = ""
    para_run(tf2, "CONFIDENTIAL  ·  PROPOSAL", SANS, 9, color=HINT, align=PP_ALIGN.RIGHT)

    # ── Main copy block ───────────────────────────────────────────────────────
    copy_y = Inches(1.5)

    # Eyebrow
    txb3 = add_textbox(slide, MARGIN_L, copy_y, USABLE_W, Inches(0.35))
    tf3 = txb3.text_frame
    tf3.paragraphs[0].text = ""
    eyebrow(tf3, "Proposal for")

    # Title
    title_y = copy_y + Inches(0.4)
    txb4 = add_textbox(slide, MARGIN_L, title_y, Inches(9), Inches(2.6))
    tf4 = txb4.text_frame
    tf4.word_wrap = True
    tf4.paragraphs[0].text = ""
    headline(tf4, "[Client Name]", size=72, italic_wrap=False, space_before_pt=0)
    headline(tf4, "A Proposal from Cognitidata", size=72, italic_wrap=True, space_before_pt=0)

    # Horizontal rule
    rule_y = title_y + Inches(2.75)
    add_rule(slide, rule_y)

    # Meta row  ·  Prepared by / Date / Project
    meta_y = rule_y + Inches(0.22)
    meta_col_w = USABLE_W / 3

    labels = ["PREPARED BY", "DATE", "ENGAGEMENT"]
    values = ["Cognitidata", "[Month Year]", "[Project Title]"]

    for i, (lbl, val) in enumerate(zip(labels, values)):
        x = MARGIN_L + meta_col_w * i
        txb_m = add_textbox(slide, x, meta_y, meta_col_w, Inches(0.6))
        tf_m = txb_m.text_frame
        tf_m.paragraphs[0].text = ""
        para_run(tf_m, lbl, SANS, 8, color=HINT, space_after_pt=3)
        para_run(tf_m, val, SANS, 12, color=TEXT, space_before_pt=2)
        # Vertical rule between columns
        if i < 2:
            add_rect(slide, x + meta_col_w - Pt(0.5), meta_y,
                     Pt(0.75), Inches(0.55), fill_color=BORDER)

    return slide


def build_outcome(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)

    # Ghost section number
    txb_n = add_textbox(slide, W - MARGIN_R - Inches(1.6), MARGIN_T - Inches(0.15),
                        Inches(1.8), Inches(1.8))
    tf_n = txb_n.text_frame
    tf_n.paragraphs[0].text = ""
    ghost_number(tf_n, "01", size=120)

    # Section header
    body_y = section_header(slide, "The Dream Outcome", "What success looks like for you", title_size=46)

    add_rule(slide, body_y + Inches(0.15))
    body_y += Inches(0.45)

    # Body placeholder
    txb_b = add_textbox(slide, MARGIN_L, body_y, Inches(7.5), Inches(3.2))
    tf_b = txb_b.text_frame
    tf_b.word_wrap = True
    tf_b.paragraphs[0].text = ""
    body_para(tf_b,
        "[Describe the client's dream outcome here. What does their life / business look like "
        "after a successful engagement? Be specific, measurable, and emotionally resonant. "
        "This paragraph should mirror the language the client used when describing their goals.]",
        size=14, color=BODY, space_before_pt=0, space_after_pt=10)
    body_para(tf_b,
        "[Secondary supporting statement — tie the outcome back to a number or milestone "
        "where possible, e.g. 'reducing manual reporting time by 80%' or 'a unified data "
        "view trusted by every team'.]",
        size=14, color=MUTED, space_before_pt=6)

    # Bottom rule + slide label
    add_rule(slide, H - Inches(0.55))
    txb_foot = add_textbox(slide, MARGIN_L, H - Inches(0.48), USABLE_W, Inches(0.3))
    tf_foot = txb_foot.text_frame
    tf_foot.paragraphs[0].text = ""
    para_run(tf_foot, "COGNITIDATA  ·  CONFIDENTIAL", SANS, 8, color=HINT)

    return slide


def build_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)

    txb_n = add_textbox(slide, W - MARGIN_R - Inches(1.6), MARGIN_T - Inches(0.15),
                        Inches(1.8), Inches(1.8))
    tf_n = txb_n.text_frame
    tf_n.paragraphs[0].text = ""
    ghost_number(tf_n, "02", size=120)

    body_y = section_header(slide, "The Problem", "What's standing in the way", title_size=46)
    add_rule(slide, body_y + Inches(0.15))
    body_y += Inches(0.45)

    problems = [
        ("01", "[Problem Title One]",
         "[Describe this specific pain point. Use the client's own words where possible. "
          "Keep it to 1–2 sentences — precise and frank.]"),
        ("02", "[Problem Title Two]",
         "[Second distinct problem. Each problem should stand alone — avoid bundling multiple "
          "issues into one item.]"),
        ("03", "[Problem Title Three]",
         "[Third problem. If fewer problems apply, delete this row.]"),
    ]

    for num, title, desc in problems:
        row_h = Inches(0.9)
        body_y = numbered_item_row(slide, num, title, desc, body_y, row_height=row_h)
        add_rule(slide, body_y, color=BORDER, thickness_pt=0.5)
        body_y += Inches(0.12)

    add_rule(slide, H - Inches(0.55))
    txb_foot = add_textbox(slide, MARGIN_L, H - Inches(0.48), USABLE_W, Inches(0.3))
    tf_foot = txb_foot.text_frame
    tf_foot.paragraphs[0].text = ""
    para_run(tf_foot, "COGNITIDATA  ·  CONFIDENTIAL", SANS, 8, color=HINT)

    return slide


def build_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)

    txb_n = add_textbox(slide, W - MARGIN_R - Inches(1.6), MARGIN_T - Inches(0.15),
                        Inches(1.8), Inches(1.8))
    tf_n = txb_n.text_frame
    tf_n.paragraphs[0].text = ""
    ghost_number(tf_n, "03", size=120)

    body_y = section_header(slide, "The Solution", "How we will get you there", title_size=46)
    add_rule(slide, body_y + Inches(0.15))
    body_y += Inches(0.45)

    solutions = [
        ("01", "[Solution Title One]",
         "[Directly addresses Problem 01. Describe the deliverable, methodology, or capability "
          "in one clear sentence. Avoid jargon unless the client speaks it.]"),
        ("02", "[Solution Title Two]",
         "[Directly addresses Problem 02. Tie each solution visually or verbally back to the "
          "corresponding problem — this creates a logical through-line.]"),
        ("03", "[Solution Title Three]",
         "[Directly addresses Problem 03. Consider ending with a short, concrete outcome: "
          ""resulting in a live dashboard by end of week 4".]"),
    ]

    for num, title, desc in solutions:
        row_h = Inches(0.9)
        body_y = numbered_item_row(slide, num, title, desc, body_y, row_height=row_h)
        add_rule(slide, body_y, color=BORDER, thickness_pt=0.5)
        body_y += Inches(0.12)

    add_rule(slide, H - Inches(0.55))
    txb_foot = add_textbox(slide, MARGIN_L, H - Inches(0.48), USABLE_W, Inches(0.3))
    tf_foot = txb_foot.text_frame
    tf_foot.paragraphs[0].text = ""
    para_run(tf_foot, "COGNITIDATA  ·  CONFIDENTIAL", SANS, 8, color=HINT)

    return slide


def build_investment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)

    txb_n = add_textbox(slide, W - MARGIN_R - Inches(1.6), MARGIN_T - Inches(0.15),
                        Inches(1.8), Inches(1.8))
    tf_n = txb_n.text_frame
    tf_n.paragraphs[0].text = ""
    ghost_number(tf_n, "04", size=120)

    body_y = section_header(slide, "The Investment", "What this engagement requires", title_size=46)
    add_rule(slide, body_y + Inches(0.15))
    body_y += Inches(0.45)

    # ── Left column: pricing ──────────────────────────────────────────────────
    col_w = (USABLE_W - Inches(0.5)) / 2

    txb_l = add_textbox(slide, MARGIN_L, body_y, col_w, Inches(3.5))
    tf_l = txb_l.text_frame
    tf_l.word_wrap = True
    tf_l.paragraphs[0].text = ""
    eyebrow(tf_l, "Investment")
    headline(tf_l, "[£ / $ Amount]", size=42, space_before_pt=6, space_after_pt=4)
    body_para(tf_l, "[Payment structure, e.g. 50% on engagement start, 50% on delivery. "
              "Or monthly retainer of £X.]", size=12, color=MUTED, space_before_pt=4)

    # Scope items
    scope_items = [
        "[Deliverable or phase 01]",
        "[Deliverable or phase 02]",
        "[Deliverable or phase 03]",
        "[Timeline: estimated X weeks]",
    ]
    for item in scope_items:
        p = tf_l.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = f"—  {item}"
        r.font.name  = SANS
        r.font.size  = Pt(12)
        r.font.color.rgb = BODY

    # ── Vertical rule between columns ─────────────────────────────────────────
    add_rect(slide,
             MARGIN_L + col_w + Inches(0.25) - Pt(0.4),
             body_y,
             Pt(0.75),
             Inches(3.4),
             fill_color=BORDER)

    # ── Right column: guarantee ───────────────────────────────────────────────
    right_x = MARGIN_L + col_w + Inches(0.5)
    txb_r = add_textbox(slide, right_x, body_y, col_w, Inches(3.5))
    tf_r = txb_r.text_frame
    tf_r.word_wrap = True
    tf_r.paragraphs[0].text = ""
    eyebrow(tf_r, "Our Guarantee")
    headline(tf_r, "You take no risk.", size=32, italic_wrap=True, space_before_pt=6, space_after_pt=4)
    body_para(tf_r,
        "[Describe your guarantee here. E.g. 'If we do not deliver the agreed outcomes within "
        "the agreed timeframe, we will continue working at no additional cost until we do.']",
        size=12, color=MUTED, space_before_pt=4)
    body_para(tf_r,
        "[Optional second guarantee or risk-reversal clause.]",
        size=12, color=MUTED, space_before_pt=8)

    add_rule(slide, H - Inches(0.55))
    txb_foot = add_textbox(slide, MARGIN_L, H - Inches(0.48), USABLE_W, Inches(0.3))
    tf_foot = txb_foot.text_frame
    tf_foot.paragraphs[0].text = ""
    para_run(tf_foot, "COGNITIDATA  ·  CONFIDENTIAL", SANS, 8, color=HINT)

    return slide


def build_next_steps(prs):
    """Simple closing slide with CTA."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)

    # Dark block left strip
    add_rect(slide, 0, 0, Inches(0.18), H, fill_color=TEXT)

    body_y = Inches(2.2)
    txb = add_textbox(slide, MARGIN_L, body_y, Inches(8), Inches(3.5))
    tf = txb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    eyebrow(tf, "Ready to begin")
    headline(tf, "Let's talk.", size=64, italic_wrap=True, space_before_pt=8, space_after_pt=8)
    body_para(tf,
        "To move forward, simply reply to this proposal or schedule a call using the link below. "
        "We'll confirm scope and get started within 48 hours.",
        size=14, color=MUTED, space_before_pt=10)
    body_para(tf, "[calendly.com/cognitidata  ·  hello@cognitidata.com]",
              size=13, color=HINT, space_before_pt=12)

    add_rule(slide, H - Inches(0.55))
    txb_foot = add_textbox(slide, MARGIN_L, H - Inches(0.48), USABLE_W, Inches(0.3))
    tf_foot = txb_foot.text_frame
    tf_foot.paragraphs[0].text = ""
    para_run(tf_foot, "COGNITIDATA  ·  CONFIDENTIAL", SANS, 8, color=HINT)

    return slide


# ── Assemble ──────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    build_cover(prs)
    build_outcome(prs)
    build_problem(prs)
    build_solution(prs)
    build_investment(prs)
    build_next_steps(prs)

    out = "cognitidata-proposal-template.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
